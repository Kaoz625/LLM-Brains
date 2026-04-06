#!/usr/bin/env python3
"""
life_data_ingest.py — Personal Life Data Ingestion
Ingests every digital trace of your life into brain/raw/ so the
LLM compiler can build a rich, interconnected personal knowledge base.

PATHWAY REPAIR THEORY:
Each data type creates a DIFFERENT KIND of retrieval pathway to the same memories.
Your iMessages contain the social context. Your GPS tracks contain the spatial
context. Your health data contains the physiological context. Your browser history
contains your attention and interest at the time. Your Spotify contains the emotional
context. When these are ALL linked to the same event in the brain, that event becomes
retrievable from any direction -- even if some pathways are damaged or forgotten.
This is the same mechanism physical therapists use to help memory-impaired patients:
build MORE routes to the same destination.

Supported data sources:
  Apple Health XML         → steps, heart rate, sleep, workouts, weight
  Google Takeout JSON      → location history, activity
  Browser history          → Chrome / Firefox / Safari (SQLite)
  iMessage chat.db         → conversations, contacts
  Email (.mbox / .eml)     → sent/received messages
  Calendar (.ics)          → events, appointments
  Contacts (.vcf)          → people directory
  Spotify history          → listening patterns, moods
  Twitter/X archive        → posts, timeline
  Documents (.docx/.xlsx/.pptx/.epub) → content extraction
  Code files               → structure and purpose (not full source)
  GPS tracks (.gpx)        → where you went

Usage:
    python life_data_ingest.py --apple-health ~/Downloads/export.xml
    python life_data_ingest.py --browser chrome
    python life_data_ingest.py --spotify ~/Downloads/my_spotify_data/
    python life_data_ingest.py --all ~/Downloads/my_data_export/
"""

import os, re, sys, json, sqlite3, logging, argparse, email, mailbox
from pathlib import Path
from datetime import datetime, timezone
from typing import Optional

log = logging.getLogger("life_data")
logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s  %(levelname)-8s  %(message)s",
                    handlers=[logging.StreamHandler(sys.stdout)])

ROOT = Path(__file__).parent / "brain"
RAW  = ROOT / "raw"

def _write_raw(name: str, content: str) -> Path:
    RAW.mkdir(parents=True, exist_ok=True)
    ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = RAW / f"{name}_{ts}.txt"
    path.write_text(content)
    log.info("Wrote: %s (%d chars)", path.name, len(content))
    return path

# ---------------------------------------------------------------------------
# Apple Health
# ---------------------------------------------------------------------------

def ingest_apple_health(xml_path: Path):
    """Parse Apple Health export.xml and write summaries to raw/."""
    try:
        import xml.etree.ElementTree as ET
    except ImportError:
        log.error("xml.etree.ElementTree not available")
        return

    log.info("Parsing Apple Health export (this may take a minute for large files)...")
    tree = ET.parse(str(xml_path))
    root = tree.getroot()

    records = {}
    for rec in root.findall("Record"):
        rtype = rec.get("type", "").replace("HKQuantityTypeIdentifier", "").replace("HKCategoryTypeIdentifier", "")
        date  = rec.get("startDate", "")[:10]
        value = rec.get("value", "")
        unit  = rec.get("unit", "")
        records.setdefault(rtype, []).append({"date": date, "value": value, "unit": unit})

    summary_lines = [
        "Source: Apple Health Export",
        f"Export date: {datetime.now().strftime('%Y-%m-%d')}",
        f"Total record types: {len(records)}",
        ""
    ]
    for rtype, entries in sorted(records.items()):
        summary_lines.append(f"## {rtype} ({len(entries)} records)")
        for e in entries[-5:]:  # last 5 of each type
            summary_lines.append(f"  {e['date']}: {e['value']} {e['unit']}")
        summary_lines.append("")

    # Workouts
    for wo in root.findall("Workout"):
        wtype    = wo.get("workoutActivityType", "").replace("HKWorkoutActivityType", "")
        duration = wo.get("duration", "?")
        date     = wo.get("startDate", "")[:10]
        summary_lines.append(f"Workout: {date} {wtype} {duration}min")

    _write_raw("apple_health", "\n".join(summary_lines))
    log.info("Apple Health: %d record types processed", len(records))

# ---------------------------------------------------------------------------
# Google Takeout Location
# ---------------------------------------------------------------------------

def ingest_google_takeout(folder: Path):
    """Process Google Takeout data folder."""
    folder = Path(folder)
    processed = 0

    # Location history
    for loc_file in folder.rglob("Records.json"):
        log.info("Processing location history: %s", loc_file)
        try:
            data = json.loads(loc_file.read_text(errors="replace"))
            locations = data.get("locations", [])
            lines = [f"Source: Google Location History\nTotal points: {len(locations)}\n"]
            for loc in locations[-200:]:  # last 200 points
                ts  = loc.get("timestamp", "")[:10]
                lat = loc.get("latitudeE7", 0) / 1e7
                lon = loc.get("longitudeE7", 0) / 1e7
                lines.append(f"{ts}: ({lat:.4f}, {lon:.4f})")
            _write_raw("google_location", "\n".join(lines))
            processed += 1
        except Exception as e:
            log.error("Location parse failed: %s", e)

    if processed == 0:
        log.warning("No recognized Google Takeout files found in %s", folder)

# ---------------------------------------------------------------------------
# Browser history
# ---------------------------------------------------------------------------

def ingest_browser_history(browser: str):
    """Extract browser history from Chrome, Firefox, or Safari SQLite DBs."""
    home = Path.home()

    if browser == "chrome":
        candidates = [
            home / "Library/Application Support/Google/Chrome/Default/History",
            home / ".config/google-chrome/Default/History",
            home / "AppData/Local/Google/Chrome/User Data/Default/History",
        ]
    elif browser == "firefox":
        ff_dir = home / "Library/Application Support/Firefox/Profiles"
        candidates = list(ff_dir.glob("*/places.sqlite")) if ff_dir.exists() else []
    elif browser == "safari":
        candidates = [home / "Library/Safari/History.db"]
    else:
        log.error("Unknown browser: %s. Use chrome, firefox, or safari", browser)
        return

    import shutil, tempfile
    for db_path in candidates:
        if not db_path.exists():
            continue
        log.info("Reading browser history: %s", db_path)

        # Copy DB (browser may have it locked)
        with tempfile.NamedTemporaryFile(suffix=".db", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        shutil.copy2(db_path, tmp_path)

        try:
            conn = sqlite3.connect(str(tmp_path))
            conn.row_factory = sqlite3.Row

            if browser in ("chrome",):
                rows = conn.execute("""
                    SELECT url, title, visit_count,
                           datetime(last_visit_time/1000000-11644473600, 'unixepoch') AS last_visit
                    FROM urls
                    ORDER BY last_visit_time DESC LIMIT 2000
                """).fetchall()
            elif browser == "firefox":
                rows = conn.execute("""
                    SELECT url, title, visit_count,
                           datetime(last_visit_date/1000000, 'unixepoch') AS last_visit
                    FROM moz_places WHERE visit_count > 0
                    ORDER BY last_visit_date DESC LIMIT 2000
                """).fetchall()
            elif browser == "safari":
                rows = conn.execute("""
                    SELECT hi.url, hv.title,
                           datetime(hv.visit_time + 978307200, 'unixepoch') AS last_visit
                    FROM history_visits hv
                    JOIN history_items hi ON hv.history_item = hi.id
                    ORDER BY hv.visit_time DESC LIMIT 2000
                """).fetchall()

            lines = [f"Source: {browser.title()} Browser History\nExtracted: {datetime.now().strftime('%Y-%m-%d')}\n"]
            for r in rows:
                lines.append(f"{r['last_visit'] or '?'}: {r.get('title','') or ''} | {r['url'][:100]}")

            _write_raw(f"browser_{browser}", "\n".join(lines))
            conn.close()
        except Exception as e:
            log.error("Browser history read failed: %s", e)
        finally:
            tmp_path.unlink(missing_ok=True)
        return

    log.warning("No %s history database found", browser)

# ---------------------------------------------------------------------------
# iMessage
# ---------------------------------------------------------------------------

def ingest_imessage(chat_db: Path):
    """Extract iMessage conversations from iPhone backup chat.db."""
    try:
        conn = sqlite3.connect(str(chat_db))
        conn.row_factory = sqlite3.Row
        rows = conn.execute("""
            SELECT m.text, m.is_from_me,
                   datetime(m.date/1000000000 + 978307200, 'unixepoch') AS msg_date,
                   h.id AS contact
            FROM message m
            LEFT JOIN handle h ON m.handle_id = h.rowid
            WHERE m.text IS NOT NULL AND m.text != ''
            ORDER BY m.date DESC LIMIT 5000
        """).fetchall()

        lines = [f"Source: iMessage\nExtracted: {datetime.now().strftime('%Y-%m-%d')}\n"]
        for r in rows:
            direction = "Me" if r["is_from_me"] else (r["contact"] or "?")
            lines.append(f"[{r['msg_date']}] {direction}: {r['text'][:200]}")

        _write_raw("imessage", "\n".join(lines))
        log.info("iMessage: %d messages extracted", len(rows))
    except Exception as e:
        log.error("iMessage read failed: %s", e)

# ---------------------------------------------------------------------------
# Email
# ---------------------------------------------------------------------------

def ingest_email(path: Path):
    """Process .mbox file or folder of .eml files."""
    path = Path(path)
    lines = [f"Source: Email\nExtracted: {datetime.now().strftime('%Y-%m-%d')}\n"]
    count = 0

    def _parse_msg(msg):
        subject = msg.get("Subject", "")
        from_   = msg.get("From", "")
        date    = msg.get("Date", "")
        body    = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    try:
                        body = part.get_payload(decode=True).decode("utf-8", errors="replace")[:500]
                    except Exception:
                        pass
                    break
        else:
            try:
                body = msg.get_payload(decode=True).decode("utf-8", errors="replace")[:500]
            except Exception:
                pass
        return f"[{date}] From: {from_} | Subject: {subject}\n{body}\n"

    if path.is_file() and path.suffix == ".mbox":
        mbox = mailbox.mbox(str(path))
        for msg in mbox:
            lines.append(_parse_msg(msg))
            count += 1
    elif path.is_dir():
        for eml in path.glob("*.eml"):
            with eml.open("rb") as f:
                msg = email.message_from_binary_file(f)
                lines.append(_parse_msg(msg))
                count += 1

    _write_raw("email", "\n".join(lines))
    log.info("Email: %d messages extracted", count)

# ---------------------------------------------------------------------------
# Calendar
# ---------------------------------------------------------------------------

def ingest_calendar(ics_path: Path):
    """Parse .ics calendar file into text entries."""
    content = Path(ics_path).read_text(errors="replace")
    lines   = [f"Source: Calendar ({ics_path.name})\n"]
    events  = []

    current = {}
    for line in content.splitlines():
        if line.startswith("BEGIN:VEVENT"):
            current = {}
        elif line.startswith("END:VEVENT"):
            events.append(current.copy())
        elif ":" in line:
            key, _, val = line.partition(":")
            key = key.split(";")[0]
            current[key] = val.strip()

    for e in events:
        summary  = e.get("SUMMARY", "")
        dtstart  = e.get("DTSTART", "")[:8]
        location = e.get("LOCATION", "")
        desc     = e.get("DESCRIPTION", "")[:200]
        lines.append(f"Event: {dtstart} | {summary} | {location}\n{desc}\n")

    _write_raw("calendar", "\n".join(lines))
    log.info("Calendar: %d events extracted", len(events))

# ---------------------------------------------------------------------------
# Contacts
# ---------------------------------------------------------------------------

def ingest_contacts(vcf_path: Path):
    """Parse .vcf vCard file."""
    content = Path(vcf_path).read_text(errors="replace")
    lines   = [f"Source: Contacts ({vcf_path.name})\n"]
    count   = 0

    card = {}
    for line in content.splitlines():
        if line.startswith("BEGIN:VCARD"):
            card = {}
        elif line.startswith("END:VCARD"):
            name   = card.get("FN", card.get("N", "Unknown"))
            phone  = card.get("TEL", "")
            email_ = card.get("EMAIL", "")
            org    = card.get("ORG", "")
            lines.append(f"Contact: {name} | {phone} | {email_} | {org}")
            count += 1
        elif ":" in line:
            key, _, val = line.partition(":")
            key = key.split(";")[0]
            card[key] = val.strip()

    _write_raw("contacts", "\n".join(lines))
    log.info("Contacts: %d cards extracted", count)

# ---------------------------------------------------------------------------
# Spotify
# ---------------------------------------------------------------------------

def ingest_spotify(folder: Path):
    """Process Spotify StreamingHistory*.json files."""
    folder = Path(folder)
    lines  = [f"Source: Spotify Listening History\n"]
    count  = 0

    for hist_file in sorted(folder.glob("StreamingHistory*.json")):
        try:
            data = json.loads(hist_file.read_text(errors="replace"))
            for entry in data:
                ts     = entry.get("endTime", "")
                artist = entry.get("artistName", "")
                track  = entry.get("trackName", "")
                ms     = entry.get("msPlayed", 0)
                if ms > 30000:  # only songs played >30s
                    lines.append(f"{ts}: {artist} — {track} ({ms//1000}s)")
                    count += 1
        except Exception as e:
            log.error("Spotify parse error %s: %s", hist_file.name, e)

    _write_raw("spotify", "\n".join(lines))
    log.info("Spotify: %d tracks extracted", count)

# ---------------------------------------------------------------------------
# Twitter/X archive
# ---------------------------------------------------------------------------

def ingest_twitter(folder: Path):
    """Process Twitter/X data export archive."""
    folder   = Path(folder)
    tweets_f = folder / "data" / "tweets.js"
    if not tweets_f.exists():
        tweets_f = folder / "tweets.js"
    if not tweets_f.exists():
        log.error("tweets.js not found in %s", folder)
        return

    # Twitter JS starts with "window.YTD.tweets.part0 = "
    raw = tweets_f.read_text(errors="replace")
    raw = re.sub(r"^window\.[^=]+=\s*", "", raw).strip()
    try:
        data   = json.loads(raw)
        tweets = [t.get("tweet", t) for t in data]
        lines  = [f"Source: Twitter/X Archive\n"]
        for t in tweets:
            created = t.get("created_at", "")
            text    = t.get("full_text", t.get("text", ""))
            likes   = t.get("favorite_count", "0")
            rts     = t.get("retweet_count", "0")
            lines.append(f"[{created}] {text[:280]}  ❤{likes} 🔁{rts}")
        _write_raw("twitter", "\n".join(lines))
        log.info("Twitter: %d tweets extracted", len(tweets))
    except Exception as e:
        log.error("Twitter parse failed: %s", e)

# ---------------------------------------------------------------------------
# Documents
# ---------------------------------------------------------------------------

def ingest_document(file_path: Path):
    """Extract text from .docx, .xlsx, .pptx, or .epub."""
    file_path = Path(file_path)
    suffix    = file_path.suffix.lower()
    text      = ""

    if suffix == ".docx":
        try:
            from docx import Document
            doc  = Document(str(file_path))
            text = "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            log.error("pip install python-docx")
            return

    elif suffix == ".xlsx":
        try:
            import openpyxl
            wb    = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append("\t".join(str(c or "") for c in row))
            text = "\n".join(lines)
        except ImportError:
            log.error("pip install openpyxl")
            return

    elif suffix == ".pptx":
        try:
            from pptx import Presentation
            prs   = Presentation(str(file_path))
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"\n--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text)
            text = "\n".join(lines)
        except ImportError:
            log.error("pip install python-pptx")
            return

    elif suffix == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from html.parser import HTMLParser

            class _Strip(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts = []
                def handle_data(self, d):
                    self.parts.append(d)
                def get_text(self):
                    return " ".join(self.parts)

            book  = epub.read_epub(str(file_path))
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                p = _Strip()
                p.feed(item.get_content().decode("utf-8", errors="replace"))
                parts.append(p.get_text())
            text = "\n".join(parts)
        except ImportError:
            log.error("pip install ebooklib")
            return

    else:
        log.error("Unsupported document type: %s", suffix)
        return

    if text.strip():
        _write_raw(f"doc_{file_path.stem[:30]}", f"Source: {file_path}\n\n{text[:20000]}")
        log.info("Document: %d chars extracted from %s", len(text), file_path.name)

# ---------------------------------------------------------------------------
# Code files
# ---------------------------------------------------------------------------

def ingest_code(file_path: Path):
    """Extract structure from code files (not full source — just purpose/shape)."""
    file_path = Path(file_path)
    source    = file_path.read_text(errors="replace")

    # Extract functions, classes, imports, top-level docstring
    lines      = source.splitlines()
    imports    = [l.strip() for l in lines if l.startswith(("import ", "from "))][:20]
    funcs      = [l.strip() for l in lines if re.match(r"^(def |async def |func |function )", l)][:30]
    classes    = [l.strip() for l in lines if re.match(r"^class ", l)][:20]
    docstring  = ""
    if source.startswith('"""') or source.startswith("'''"):
        end = source.find('"""', 3) if source.startswith('"""') else source.find("'''", 3)
        if end > 0:
            docstring = source[3:end].strip()[:500]

    content = (
        f"Source: {file_path}\nLanguage: {file_path.suffix}\nLines: {len(lines)}\n\n"
        f"Purpose:\n{docstring}\n\n"
        f"Imports:\n" + "\n".join(imports) + "\n\n"
        f"Classes:\n" + "\n".join(classes) + "\n\n"
        f"Functions:\n" + "\n".join(funcs)
    )
    _write_raw(f"code_{file_path.stem[:30]}", content)

# ---------------------------------------------------------------------------
# GPS tracks
# ---------------------------------------------------------------------------

def ingest_gpx(gpx_path: Path):
    """Parse .gpx GPS track file."""
    import xml.etree.ElementTree as ET
    try:
        tree   = ET.parse(str(gpx_path))
        root   = tree.getroot()
        ns     = {"gpx": "http://www.topografix.com/GPX/1/1"}
        points = root.findall(".//gpx:trkpt", ns) or root.findall(".//trkpt")

        lines = [f"Source: GPS Track ({gpx_path.name})\nPoints: {len(points)}\n"]
        for pt in points[::max(1, len(points)//100)]:  # sample max 100 points
            lat  = pt.get("lat", "")
            lon  = pt.get("lon", "")
            time = pt.findtext("{http://www.topografix.com/GPX/1/1}time", pt.findtext("time", ""))
            ele  = pt.findtext("{http://www.topografix.com/GPX/1/1}ele", pt.findtext("ele", ""))
            lines.append(f"{time[:16] if time else '?'}: ({lat}, {lon}) ele={ele}m")

        _write_raw(f"gps_{gpx_path.stem[:30]}", "\n".join(lines))
        log.info("GPX: %d track points extracted", len(points))
    except Exception as e:
        log.error("GPX parse failed: %s", e)

# ---------------------------------------------------------------------------
# Auto-detect from folder (--all)
# ---------------------------------------------------------------------------

def ingest_all(folder: Path):
    """Walk a folder and try to auto-detect and process all known data types."""
    folder = Path(folder)
    log.info("Auto-detecting data types in %s", folder)

    for f in folder.rglob("export.xml"):
        if f.stat().st_size > 10000:
            log.info("Found Apple Health export: %s", f)
            ingest_apple_health(f)

    for f in folder.rglob("StreamingHistory*.json"):
        log.info("Found Spotify history: %s", f.parent)
        ingest_spotify(f.parent)
        break

    for f in folder.rglob("tweets.js"):
        log.info("Found Twitter archive: %s", f.parent.parent)
        ingest_twitter(f.parent.parent)

    for f in folder.rglob("Records.json"):
        if "Location" in str(f) or "location" in str(f):
            log.info("Found Google location history: %s", f.parent)
            ingest_google_takeout(f.parent)

    for f in folder.rglob("*.ics"):
        log.info("Found calendar: %s", f)
        ingest_calendar(f)

    for f in folder.rglob("*.vcf"):
        log.info("Found contacts: %s", f)
        ingest_contacts(f)

    for f in folder.rglob("*.gpx"):
        log.info("Found GPS track: %s", f)
        ingest_gpx(f)

    for f in folder.rglob("*.mbox"):
        log.info("Found email mbox: %s", f)
        ingest_email(f)

    for f in folder.rglob("*.eml"):
        log.info("Found email folder: %s", f.parent)
        ingest_email(f.parent)
        break

    for ext in (".docx", ".epub", ".pptx"):
        for f in folder.rglob(f"*{ext}"):
            ingest_document(f)
life_data_ingest.py — Personal life data ingestion from various export formats.

Supports:
  - Apple Health XML
  - Google Takeout JSON (location history)
  - iMessage export
  - .ics calendar files
  - .vcf contacts
  - Browser history (Chrome/Safari/Firefox)
  - Spotify history JSON
  - .mbox/.eml email
  - .gpx GPS tracks

Usage:
    python life_data_ingest.py --apple-health file.xml
    python life_data_ingest.py --calendar file.ics
    python life_data_ingest.py --all ~/Downloads/Export/
    python life_data_ingest.py --contacts contacts.vcf
"""

import argparse
import csv
import json
import os
import re
import sqlite3
import sys
import xml.etree.ElementTree as ET
from collections import defaultdict
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv

load_dotenv()

BRAIN_DIR = Path(os.getenv("BRAIN_DIR", "./brain"))
ME_DIR = BRAIN_DIR / "me"


def ensure_dirs():
    ME_DIR.mkdir(parents=True, exist_ok=True)
    (BRAIN_DIR / "knowledge").mkdir(parents=True, exist_ok=True)


def append_to_file(path: Path, content: str):
    """Append content to a file, creating it if needed."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "a", encoding="utf-8") as f:
        f.write(content + "\n")


def get_or_create_file(path: Path, header: str = "") -> Path:
    if not path.exists():
        path.write_text(header + "\n", encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# Apple Health XML
# ---------------------------------------------------------------------------

def parse_apple_health(xml_path: Path) -> dict:
    """Parse Apple Health export XML and extract key stats."""
    print(f"  Parsing Apple Health: {xml_path.name}")
    stats = defaultdict(list)

    try:
        tree = ET.parse(str(xml_path))
        root = tree.getroot()
    except ET.ParseError as e:
        print(f"  XML parse error: {e}", file=sys.stderr)
        return {}

    record_types_of_interest = {
        "HKQuantityTypeIdentifierHeartRate": "heart_rate",
        "HKQuantityTypeIdentifierStepCount": "steps",
        "HKQuantityTypeIdentifierBodyMass": "weight",
        "HKQuantityTypeIdentifierActiveEnergyBurned": "active_calories",
        "HKCategoryTypeIdentifierSleepAnalysis": "sleep",
        "HKQuantityTypeIdentifierVO2Max": "vo2_max",
        "HKQuantityTypeIdentifierBloodPressureSystolic": "bp_systolic",
        "HKQuantityTypeIdentifierBloodPressureDiastolic": "bp_diastolic",
        "HKQuantityTypeIdentifierDistanceWalkingRunning": "distance_km",
        "HKQuantityTypeIdentifierMindfulSession": "mindfulness",
    }

    for record in root.iter("Record"):
        rtype = record.get("type", "")
        if rtype in record_types_of_interest:
            key = record_types_of_interest[rtype]
            value = record.get("value", "")
            unit = record.get("unit", "")
            start_date = record.get("startDate", "")[:10]
            try:
                stats[key].append({
                    "date": start_date,
                    "value": float(value) if value else 0,
                    "unit": unit,
                })
            except ValueError:
                pass

    return dict(stats)


def write_health_summary(stats: dict, output_path: Path):
    """Write Apple Health stats to health.md."""
    now = datetime.now().strftime("%Y-%m-%d")
    lines = [f"\n## Health Import — {now}\n"]

    def avg_recent(data: list, days: int = 30) -> float:
        cutoff = (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d")
        recent = [d["value"] for d in data if d.get("date", "") >= cutoff]
        return sum(recent) / len(recent) if recent else 0

    if "steps" in stats:
        avg_steps = avg_recent(stats["steps"])
        lines.append(f"### Steps\n- **30-day avg:** {avg_steps:.0f} steps/day")

    if "heart_rate" in stats:
        recent_hr = [d["value"] for d in stats["heart_rate"][-50:]]
        if recent_hr:
            lines.append(f"\n### Heart Rate\n- **Recent avg:** {sum(recent_hr)/len(recent_hr):.0f} bpm")
            lines.append(f"- **Min:** {min(recent_hr):.0f} | **Max:** {max(recent_hr):.0f}")

    if "weight" in stats:
        recent_weights = sorted(stats["weight"], key=lambda x: x["date"])[-5:]
        if recent_weights:
            latest = recent_weights[-1]
            lines.append(f"\n### Weight\n- **Latest:** {latest['value']:.1f} {latest['unit']} ({latest['date']})")

    if "sleep" in stats:
        sleep_records = [d for d in stats["sleep"] if d["value"] > 0]
        lines.append(f"\n### Sleep\n- **Records:** {len(sleep_records)} nights tracked")

    if "vo2_max" in stats:
        recent_vo2 = sorted(stats["vo2_max"], key=lambda x: x["date"])[-1:]
        if recent_vo2:
            lines.append(f"\n### VO2 Max\n- **Latest:** {recent_vo2[0]['value']:.1f} mL/kg/min")

    content = "\n".join(lines)
    append_to_file(output_path, content)
    print(f"  -> Health data written to {output_path.name}")


# ---------------------------------------------------------------------------
# Google Location History JSON
# ---------------------------------------------------------------------------

def parse_google_location(json_path: Path) -> list[dict]:
    """Parse Google Takeout location history."""
    print(f"  Parsing Google Location: {json_path.name}")
    locations = []

    try:
        data = json.loads(json_path.read_text(encoding="utf-8", errors="replace"))
        # Handle both old and new Takeout formats
        records = data.get("locations", data.get("timelineObjects", []))

        for record in records[:1000]:  # Limit to first 1000
            if isinstance(record, dict):
                # Old format
                if "latitudeE7" in record:
                    lat = record["latitudeE7"] / 1e7
                    lng = record["longitudeE7"] / 1e7
                    ts = int(record.get("timestampMs", 0)) / 1000
                    dt = datetime.fromtimestamp(ts).strftime("%Y-%m-%d") if ts else ""
                    locations.append({"date": dt, "lat": lat, "lng": lng})
                # New format
                elif "placeVisit" in record:
                    visit = record["placeVisit"]
                    loc = visit.get("location", {})
                    locations.append({
                        "date": visit.get("duration", {}).get("startTimestampMs", "")[:10],
                        "name": loc.get("name", ""),
                        "address": loc.get("address", ""),
                        "lat": loc.get("latitudeE7", 0) / 1e7,
                        "lng": loc.get("longitudeE7", 0) / 1e7,
                    })
    except Exception as e:
        print(f"  Location parse error: {e}", file=sys.stderr)

    return locations


def write_locations_summary(locations: list[dict], output_path: Path):
    """Write location data to locations.md."""
    now = datetime.now().strftime("%Y-%m-%d")
    lines = [f"\n## Location Import — {now}\n"]
    lines.append(f"**Total records:** {len(locations)}\n")

    # Group by named places
    named = [l for l in locations if l.get("name")]
    if named:
        place_counts = defaultdict(int)
        for l in named:
            place_counts[l["name"]] += 1
        top_places = sorted(place_counts.items(), key=lambda x: x[1], reverse=True)[:20]
        lines.append("\n### Most Visited Places")
        for place, count in top_places:
            lines.append(f"- **{place}**: {count} visits")

    append_to_file(output_path, "\n".join(lines))
    print(f"  -> {len(locations)} locations written to {output_path.name}")


# ---------------------------------------------------------------------------
# Calendar ICS
# ---------------------------------------------------------------------------

def parse_ics(ics_path: Path) -> list[dict]:
    """Parse .ics calendar file and extract events."""
    print(f"  Parsing calendar: {ics_path.name}")
    events = []

    try:
        content = ics_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        print(f"  ICS read error: {e}", file=sys.stderr)
        return []

    # Simple ICS parser (without icalendar library)
    current_event = {}
    in_event = False

    for line in content.splitlines():
        line = line.strip()
        if line == "BEGIN:VEVENT":
            in_event = True
            current_event = {}
        elif line == "END:VEVENT" and in_event:
            if current_event.get("SUMMARY"):
                events.append(current_event.copy())
            in_event = False
            current_event = {}
        elif in_event and ":" in line:
            key, _, value = line.partition(":")
            # Handle property parameters (e.g., DTSTART;TZID=...)
            key = key.split(";")[0]
            current_event[key] = value

    return events


def format_ics_date(dt_str: str) -> str:
    """Parse ICS datetime string to readable format."""
    dt_str = dt_str.strip()
    for fmt in ("%Y%m%dT%H%M%SZ", "%Y%m%dT%H%M%S", "%Y%m%d"):
        try:
            return datetime.strptime(dt_str, fmt).strftime("%Y-%m-%d %H:%M")
        except ValueError:
            pass
    return dt_str


def write_calendar_events(events: list[dict], output_path: Path):
    """Write calendar events to timeline.md."""
    now = datetime.now().strftime("%Y-%m-%d")
    lines = [f"\n## Calendar Import — {now} ({len(events)} events)\n"]

    # Sort by date
    sorted_events = sorted(events, key=lambda e: e.get("DTSTART", ""), reverse=True)

    for event in sorted_events[:100]:  # Limit display
        summary = event.get("SUMMARY", "Untitled")
        start = format_ics_date(event.get("DTSTART", ""))
        end = format_ics_date(event.get("DTEND", ""))
        location = event.get("LOCATION", "")
        description = event.get("DESCRIPTION", "")[:100]

        lines.append(f"### {start} — {summary}")
        if location:
            lines.append(f"**Location:** {location}")
        if end and end != start:
            lines.append(f"**End:** {end}")
        if description:
            lines.append(f"*{description}*")
        lines.append("")

    append_to_file(output_path, "\n".join(lines))
    print(f"  -> {len(events)} events written to {output_path.name}")


# ---------------------------------------------------------------------------
# VCF Contacts
# ---------------------------------------------------------------------------

def parse_vcf(vcf_path: Path) -> list[dict]:
    """Parse .vcf contacts file."""
    print(f"  Parsing contacts: {vcf_path.name}")
    contacts = []
    current = {}

    try:
        content = vcf_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        print(f"  VCF read error: {e}", file=sys.stderr)
        return []

    for line in content.splitlines():
        line = line.strip()
        if line == "BEGIN:VCARD":
            current = {}
        elif line == "END:VCARD":
            if current.get("FN") or current.get("N"):
                contacts.append(current.copy())
            current = {}
        elif ":" in line:
            key, _, value = line.partition(":")
            key = key.split(";")[0]
            if key in ("FN", "N", "EMAIL", "TEL", "ORG", "TITLE", "URL", "NOTE", "BDAY"):
                current[key] = value

    return contacts


def write_contacts(contacts: list[dict], output_dir: Path):
    """Write contacts as wiki entries in knowledge/people/."""
    people_dir = output_dir / "knowledge" / "people"
    people_dir.mkdir(parents=True, exist_ok=True)

    written = 0
    for contact in contacts:
        name = contact.get("FN", "")
        if not name:
            name_parts = contact.get("N", "").split(";")
            name = " ".join(p for p in reversed(name_parts[:2]) if p).strip()
        if not name:
            continue

        slug = re.sub(r"[^\w-]", "-", name.lower())[:40]
        md_path = people_dir / f"{slug}.md"

        content = f"""---
title: {name}
type: person
email: {contact.get('EMAIL', '')}
phone: {contact.get('TEL', '')}
org: {contact.get('ORG', '')}
title_role: {contact.get('TITLE', '')}
birthday: {contact.get('BDAY', '')}
imported: {datetime.now().isoformat()}
---

# {name}

**Organization:** {contact.get('ORG', 'Unknown')}
**Role:** {contact.get('TITLE', '')}
**Email:** {contact.get('EMAIL', '')}
**Phone:** {contact.get('TEL', '')}
**Website:** {contact.get('URL', '')}

## Notes
{contact.get('NOTE', '')}

## Relationship
<!-- Add relationship notes here -->

"""
        md_path.write_text(content, encoding="utf-8")
        written += 1

    print(f"  -> {written} contacts written to {people_dir.relative_to(BRAIN_DIR)}")


# ---------------------------------------------------------------------------
# Browser History
# ---------------------------------------------------------------------------

def parse_chrome_history(history_path: Path) -> list[dict]:
    """Parse Chrome History SQLite database."""
    print(f"  Parsing Chrome history: {history_path.name}")
    visits = []

    try:
        import shutil
        import tempfile
        # Copy DB (Chrome might have it locked)
        with tempfile.NamedTemporaryFile(suffix=".db", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        shutil.copy2(str(history_path), str(tmp_path))

        conn = sqlite3.connect(str(tmp_path))
        conn.row_factory = sqlite3.Row
        rows = conn.execute("""
            SELECT u.url, u.title, v.visit_time
            FROM urls u
            JOIN visits v ON u.id = v.url
            ORDER BY v.visit_time DESC
            LIMIT 5000
        """).fetchall()
        conn.close()
        tmp_path.unlink()

        for row in rows:
            # Chrome timestamps are microseconds since 1601-01-01
            ts = row["visit_time"]
            if ts:
                epoch_offset = 11644473600  # seconds between 1601 and 1970
                dt = datetime.fromtimestamp(ts / 1e6 - epoch_offset)
                date_str = dt.strftime("%Y-%m-%d")
            else:
                date_str = ""
            visits.append({
                "url": row["url"],
                "title": row["title"],
                "date": date_str,
            })
    except Exception as e:
        print(f"  Chrome history error: {e}", file=sys.stderr)

    return visits


def write_browser_interests(visits: list[dict], output_path: Path):
    """Extract reading interests from browser history."""
    now = datetime.now().strftime("%Y-%m-%d")
    domain_counts = defaultdict(int)
    for v in visits:
        try:
            from urllib.parse import urlparse
            domain = urlparse(v["url"]).netloc.replace("www.", "")
            domain_counts[domain] += 1
        except Exception:
            pass

    top_domains = sorted(domain_counts.items(), key=lambda x: x[1], reverse=True)[:30]
    lines = [f"\n## Browser History Import — {now}\n",
             f"**Total visits:** {len(visits)}\n",
             "\n### Top Sites"]
    for domain, count in top_domains:
        lines.append(f"- **{domain}**: {count} visits")

    append_to_file(output_path, "\n".join(lines))
    print(f"  -> Browser history written to {output_path.name}")


# ---------------------------------------------------------------------------
# Spotify History
# ---------------------------------------------------------------------------

def parse_spotify_history(json_path: Path) -> list[dict]:
    """Parse Spotify streaming history JSON."""
    print(f"  Parsing Spotify history: {json_path.name}")
    try:
        data = json.loads(json_path.read_text(encoding="utf-8", errors="replace"))
        if isinstance(data, list):
            return data
    except Exception as e:
        print(f"  Spotify parse error: {e}", file=sys.stderr)
    return []


def write_spotify_analysis(plays: list[dict], output_path: Path):
    """Write Spotify listening patterns."""
    now = datetime.now().strftime("%Y-%m-%d")
    artist_counts = defaultdict(int)
    track_counts = defaultdict(int)

    for play in plays:
        artist = play.get("artistName", play.get("master_metadata_album_artist_name", ""))
        track = play.get("trackName", play.get("master_metadata_track_name", ""))
        if artist:
            artist_counts[artist] += 1
        if track:
            track_counts[track] += 1

    top_artists = sorted(artist_counts.items(), key=lambda x: x[1], reverse=True)[:20]
    top_tracks = sorted(track_counts.items(), key=lambda x: x[1], reverse=True)[:20]

    lines = [f"\n## Spotify History Import — {now}\n",
             f"**Total plays:** {len(plays)}\n",
             "\n### Top Artists"]
    for artist, count in top_artists:
        lines.append(f"- **{artist}**: {count} plays")
    lines.append("\n### Top Tracks")
    for track, count in top_tracks:
        lines.append(f"- **{track}**: {count} plays")

    append_to_file(output_path, "\n".join(lines))
    print(f"  -> Spotify data written to {output_path.name}")


# ---------------------------------------------------------------------------
# Email (.mbox/.eml)
# ---------------------------------------------------------------------------

def parse_email(email_path: Path) -> list[dict]:
    """Parse email file(s) and extract relationship data."""
    print(f"  Parsing email: {email_path.name}")
    messages = []

    try:
        import mailbox
        if email_path.suffix.lower() == ".mbox":
            mbox = mailbox.mbox(str(email_path))
            for msg in mbox:
                messages.append({
                    "from": str(msg.get("From", "")),
                    "to": str(msg.get("To", "")),
                    "subject": str(msg.get("Subject", "")),
                    "date": str(msg.get("Date", ""))[:10],
                })
        else:
            import email as email_lib
            with open(email_path, "rb") as f:
                msg = email_lib.message_from_bytes(f.read())
            messages.append({
                "from": str(msg.get("From", "")),
                "to": str(msg.get("To", "")),
                "subject": str(msg.get("Subject", "")),
                "date": str(msg.get("Date", ""))[:10],
            })
    except Exception as e:
        print(f"  Email parse error: {e}", file=sys.stderr)

    return messages


def write_email_relationships(messages: list[dict], output_path: Path):
    """Extract relationship map from email data."""
    now = datetime.now().strftime("%Y-%m-%d")
    contact_counts = defaultdict(int)

    email_pattern = re.compile(r'[\w.+-]+@[\w-]+\.[a-zA-Z]{2,}')
    for msg in messages:
        for field in ("from", "to"):
            emails = email_pattern.findall(msg.get(field, ""))
            for e in emails:
                if "noreply" not in e and "no-reply" not in e:
                    contact_counts[e.lower()] += 1

    top_contacts = sorted(contact_counts.items(), key=lambda x: x[1], reverse=True)[:30]
    lines = [f"\n## Email Import — {now}\n",
             f"**Total messages:** {len(messages)}\n",
             "\n### Frequent Contacts"]
    for contact, count in top_contacts:
        lines.append(f"- **{contact}**: {count} messages")

    append_to_file(output_path, "\n".join(lines))
    print(f"  -> Email relationship map written to {output_path.name}")


# ---------------------------------------------------------------------------
# GPX GPS tracks
# ---------------------------------------------------------------------------

def parse_gpx(gpx_path: Path) -> list[dict]:
    """Parse GPX GPS track file."""
    print(f"  Parsing GPX: {gpx_path.name}")
    tracks = []

    try:
        tree = ET.parse(str(gpx_path))
        root = tree.getroot()
        ns = {"gpx": "http://www.topografix.com/GPX/1/1"}

        for trkpt in root.iter("{http://www.topografix.com/GPX/1/1}trkpt"):
            lat = float(trkpt.get("lat", 0))
            lon = float(trkpt.get("lon", 0))
            time_elem = trkpt.find("{http://www.topografix.com/GPX/1/1}time")
            time_str = time_elem.text[:10] if time_elem is not None and time_elem.text else ""
            ele_elem = trkpt.find("{http://www.topografix.com/GPX/1/1}ele")
            ele = float(ele_elem.text) if ele_elem is not None and ele_elem.text else 0

            tracks.append({"lat": lat, "lon": lon, "time": time_str, "elevation": ele})
    except Exception as e:
        print(f"  GPX parse error: {e}", file=sys.stderr)

    return tracks


def write_gpx_summary(tracks: list[dict], source_path: Path, output_path: Path):
    """Write GPX track summary."""
    if not tracks:
        return
    now = datetime.now().strftime("%Y-%m-%d")
    start = tracks[0]
    end = tracks[-1]

    # Estimate total distance using Haversine
    import math
    total_km = 0
    for i in range(1, len(tracks)):
        lat1, lon1 = math.radians(tracks[i-1]["lat"]), math.radians(tracks[i-1]["lon"])
        lat2, lon2 = math.radians(tracks[i]["lat"]), math.radians(tracks[i]["lon"])
        dlat = lat2 - lat1
        dlon = lon2 - lon1
        a = math.sin(dlat/2)**2 + math.cos(lat1) * math.cos(lat2) * math.sin(dlon/2)**2
        total_km += 6371 * 2 * math.asin(math.sqrt(a))

    elevations = [t["elevation"] for t in tracks if t["elevation"] > 0]
    elev_gain = max(elevations) - min(elevations) if elevations else 0

    lines = [
        f"\n## GPS Track — {source_path.name} ({now})\n",
        f"- **Points:** {len(tracks)}",
        f"- **Date:** {start.get('time', 'Unknown')} to {end.get('time', 'Unknown')}",
        f"- **Distance:** {total_km:.2f} km",
        f"- **Elevation range:** {min(elevations):.0f}m – {max(elevations):.0f}m" if elevations else "",
        f"- **Start:** {start['lat']:.4f}, {start['lon']:.4f}",
        f"- **End:** {end['lat']:.4f}, {end['lon']:.4f}",
    ]
    append_to_file(output_path, "\n".join(l for l in lines if l))
    print(f"  -> GPX summary written ({total_km:.1f} km)")


# ---------------------------------------------------------------------------
# Directory auto-detect ingest
# ---------------------------------------------------------------------------

def ingest_directory(directory: Path):
    """Auto-detect and ingest all supported files in a directory."""
    print(f"\nAuto-ingesting directory: {directory}")
    ensure_dirs()

    for path in sorted(directory.rglob("*")):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        name = path.name.lower()

        try:
            if suffix == ".xml" and ("health" in name or "export" in name):
                stats = parse_apple_health(path)
                if stats:
                    write_health_summary(stats, ME_DIR / "health.md")

            elif suffix == ".json" and "location" in name:
                locs = parse_google_location(path)
                if locs:
                    write_locations_summary(locs, ME_DIR / "locations.md")

            elif suffix == ".json" and "spotify" in name.lower():
                plays = parse_spotify_history(path)
                if plays:
                    write_spotify_analysis(plays, ME_DIR / "music.md")

            elif suffix == ".ics":
                events = parse_ics(path)
                if events:
                    write_calendar_events(events, ME_DIR / "timeline.md")

            elif suffix == ".vcf":
                contacts = parse_vcf(path)
                if contacts:
                    write_contacts(contacts, BRAIN_DIR)

            elif name in ("history",) or (suffix == ".db" and "history" in name):
                visits = parse_chrome_history(path)
                if visits:
                    write_browser_interests(visits, ME_DIR / "interests.md")

            elif suffix in (".mbox", ".eml"):
                msgs = parse_email(path)
                if msgs:
                    write_email_relationships(msgs, ME_DIR / "relationships.md")

            elif suffix == ".gpx":
                tracks = parse_gpx(path)
                if tracks:
                    write_gpx_summary(tracks, path, ME_DIR / "locations.md")

        except Exception as e:
            print(f"  Error processing {path.name}: {e}", file=sys.stderr)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Life data ingestion — turns your digital life into brain pathways"
    )
    parser.add_argument("--apple-health",  metavar="PATH", help="Apple Health export.xml")
    parser.add_argument("--google-takeout",metavar="PATH", help="Google Takeout folder")
    parser.add_argument("--browser",       metavar="NAME", help="chrome | firefox | safari")
    parser.add_argument("--imessage",      metavar="PATH", help="iPhone backup chat.db")
    parser.add_argument("--email",         metavar="PATH", help=".mbox file or .eml folder")
    parser.add_argument("--calendar",      metavar="PATH", help=".ics calendar file")
    parser.add_argument("--contacts",      metavar="PATH", help=".vcf contacts file")
    parser.add_argument("--spotify",       metavar="PATH", help="Spotify data download folder")
    parser.add_argument("--twitter",       metavar="PATH", help="Twitter archive folder")
    parser.add_argument("--document",      metavar="PATH", help=".docx/.xlsx/.pptx/.epub file")
    parser.add_argument("--code",          metavar="PATH", help="Code file to index structure")
    parser.add_argument("--gpx",           metavar="PATH", help=".gpx GPS track file")
    parser.add_argument("--all",           metavar="PATH", help="Auto-detect everything in folder")
    args = parser.parse_args()

    ROOT.mkdir(parents=True, exist_ok=True)

    if args.apple_health:
        ingest_apple_health(Path(args.apple_health))
    if args.google_takeout:
        ingest_google_takeout(Path(args.google_takeout))
    if args.browser:
        ingest_browser_history(args.browser)
    if args.imessage:
        ingest_imessage(Path(args.imessage))
    if args.email:
        ingest_email(Path(args.email))
    if args.calendar:
        ingest_calendar(Path(args.calendar))
    if args.contacts:
        ingest_contacts(Path(args.contacts))
    if args.spotify:
        ingest_spotify(Path(args.spotify))
    if args.twitter:
        ingest_twitter(Path(args.twitter))
    if args.document:
        ingest_document(Path(args.document))
    if args.code:
        ingest_code(Path(args.code))
    if args.gpx:
        ingest_gpx(Path(args.gpx))
    if args.all:
        ingest_all(Path(args.all))

    if not any(vars(args).values()):
        parser.print_help()
        print("\nAfter ingesting, run: python compile.py")
    parser = argparse.ArgumentParser(description="Personal life data ingestion")
    parser.add_argument("--apple-health", type=str, metavar="FILE",
                        help="Apple Health export.xml")
    parser.add_argument("--location", type=str, metavar="FILE",
                        help="Google Takeout location history JSON")
    parser.add_argument("--calendar", type=str, metavar="FILE",
                        help=".ics calendar file")
    parser.add_argument("--contacts", type=str, metavar="FILE",
                        help=".vcf contacts file")
    parser.add_argument("--spotify", type=str, metavar="FILE",
                        help="Spotify streaming history JSON")
    parser.add_argument("--email", type=str, metavar="FILE",
                        help=".mbox or .eml email file")
    parser.add_argument("--gpx", type=str, metavar="FILE",
                        help=".gpx GPS track file")
    parser.add_argument("--browser-history", type=str, metavar="FILE",
                        help="Chrome History database file")
    parser.add_argument("--all", type=str, metavar="DIR",
                        help="Auto-detect and ingest all files in directory")
    args = parser.parse_args()

    ensure_dirs()

    if args.all:
        ingest_directory(Path(args.all))

    else:
        processed = False
        if args.apple_health:
            stats = parse_apple_health(Path(args.apple_health))
            write_health_summary(stats, ME_DIR / "health.md")
            processed = True

        if args.location:
            locs = parse_google_location(Path(args.location))
            write_locations_summary(locs, ME_DIR / "locations.md")
            processed = True

        if args.calendar:
            events = parse_ics(Path(args.calendar))
            write_calendar_events(events, ME_DIR / "timeline.md")
            processed = True

        if args.contacts:
            contacts = parse_vcf(Path(args.contacts))
            write_contacts(contacts, BRAIN_DIR)
            processed = True

        if args.spotify:
            plays = parse_spotify_history(Path(args.spotify))
            write_spotify_analysis(plays, ME_DIR / "music.md")
            processed = True

        if args.email:
            msgs = parse_email(Path(args.email))
            write_email_relationships(msgs, ME_DIR / "relationships.md")
            processed = True

        if args.gpx:
            tracks = parse_gpx(Path(args.gpx))
            write_gpx_summary(tracks, Path(args.gpx), ME_DIR / "locations.md")
            processed = True

        if args.browser_history:
            visits = parse_chrome_history(Path(args.browser_history))
            write_browser_interests(visits, ME_DIR / "interests.md")
            processed = True

        if not processed:
            parser.print_help()


if __name__ == "__main__":
    main()
